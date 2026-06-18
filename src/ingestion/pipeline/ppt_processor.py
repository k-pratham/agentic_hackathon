"""
PPTX slide processing utilities.

Uses python-pptx to extract text content, bullet points, and tables
from presentation slides into markdown representations.
"""

import os
import logging
from pptx import Presentation

logger = logging.getLogger(__name__)


def parse_pptx_to_markdown_slides(file_path: str) -> list[str]:
    """
    Parse a PPTX file and return a list of markdown representations (one per slide).

    For each slide:
      - Extracts the slide title (if present)
      - Extracts paragraph text, bullet points, and shapes
      - Formats tables inside slides as markdown tables

    Args:
      file_path: Absolute path to the PPTX file.

    Returns:
      List of strings, each containing the formatted content of a slide.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PPTX file not found: {file_path}")

    logger.info("Parsing PPTX presentation: %s", os.path.basename(file_path))
    prs = Presentation(file_path)
    slides_content: list[str] = []

    for index, slide in enumerate(prs.slides):
        slide_num = index + 1
        slide_lines = [f"## Slide {slide_num}"]

        # 1. Try to find slide title
        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                slide_lines.append(f"# {title_text}")

        # 2. Iterate through shapes to extract text and tables
        for shape in slide.shapes:
            # Skip title shape as it was already handled
            if slide.shapes.title and shape == slide.shapes.title:
                continue

            # Handle text frame (bullet points, paragraphs)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    # Format as bullet if indented or has bullet style
                    indent = paragraph.level
                    if indent > 0:
                        slide_lines.append("  " * indent + f"- {text}")
                    else:
                        slide_lines.append(text)

            # Handle table shapes
            if shape.has_table:
                table = shape.table
                table_lines = []
                # Simple markdown table generation
                headers = []
                for col_idx in range(len(table.columns)):
                    cell = table.cell(0, col_idx)
                    headers.append(cell.text.strip() or " ")
                
                table_lines.append("| " + " | ".join(headers) + " |")
                table_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

                for row_idx in range(1, len(table.rows)):
                    row_cells = []
                    for col_idx in range(len(table.columns)):
                        cell = table.cell(row_idx, col_idx)
                        row_cells.append(cell.text.replace("\n", " ").strip() or " ")
                    table_lines.append("| " + " | ".join(row_cells) + " |")

                slide_lines.append("\n" + "\n".join(table_lines) + "\n")

        # Join lines for the slide
        slide_text = "\n".join(slide_lines).strip()
        if slide_text:
            slides_content.append(slide_text)

    logger.info(
        "Successfully parsed %d slides from %s",
        len(slides_content),
        os.path.basename(file_path),
    )
    return slides_content
